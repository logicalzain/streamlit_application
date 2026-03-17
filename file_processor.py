"""
File Processor Module
Handles reading and extracting text from various file formats:
PDF, TXT, DOCX, XLSX, CSV, PPTX, JSON, HTML, Markdown
"""

import io
import json
import csv
import os


def extract_text_from_file(uploaded_file) -> str:
    """Extract text content from an uploaded file based on its extension."""
    filename = uploaded_file.name.lower()
    raw_bytes = uploaded_file.read()
    uploaded_file.seek(0)

    if filename.endswith(".pdf"):
        return _extract_pdf(raw_bytes)
    elif filename.endswith(".txt") or filename.endswith(".md"):
        return raw_bytes.decode("utf-8", errors="replace")
    elif filename.endswith(".docx"):
        return _extract_docx(raw_bytes)
    elif filename.endswith(".xlsx") or filename.endswith(".xls"):
        return _extract_excel(raw_bytes)
    elif filename.endswith(".csv"):
        return _extract_csv(raw_bytes)
    elif filename.endswith(".pptx"):
        return _extract_pptx(raw_bytes)
    elif filename.endswith(".json"):
        return _extract_json(raw_bytes)
    elif filename.endswith(".html") or filename.endswith(".htm"):
        return _extract_html(raw_bytes)
    else:
        # Try plain text as fallback
        try:
            return raw_bytes.decode("utf-8", errors="replace")
        except Exception:
            return "[Could not extract text from this file format]"


def _extract_pdf(raw_bytes: bytes) -> str:
    """Extract text from PDF bytes using PyPDF2."""
    from PyPDF2 import PdfReader
    reader = PdfReader(io.BytesIO(raw_bytes))
    pages = []
    for i, page in enumerate(reader.pages):
        text = page.extract_text() or ""
        if text.strip():
            pages.append(f"--- Page {i + 1} ---\n{text}")
    return "\n\n".join(pages) if pages else "[No extractable text found in PDF]"


def _extract_docx(raw_bytes: bytes) -> str:
    """Extract text from DOCX bytes using python-docx."""
    from docx import Document
    doc = Document(io.BytesIO(raw_bytes))
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
    # Also extract text from tables
    for table in doc.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if cells:
                paragraphs.append(" | ".join(cells))
    return "\n\n".join(paragraphs) if paragraphs else "[No text found in document]"


def _extract_excel(raw_bytes: bytes) -> str:
    """Extract text from Excel bytes using openpyxl."""
    from openpyxl import load_workbook
    wb = load_workbook(io.BytesIO(raw_bytes), read_only=True, data_only=True)
    all_text = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        all_text.append(f"=== Sheet: {sheet_name} ===")
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) if c is not None else "" for c in row]
            if any(c.strip() for c in cells):
                all_text.append(" | ".join(cells))
    wb.close()
    return "\n".join(all_text) if all_text else "[No data found in spreadsheet]"


def _extract_csv(raw_bytes: bytes) -> str:
    """Extract text from CSV bytes."""
    text = raw_bytes.decode("utf-8", errors="replace")
    reader = csv.reader(io.StringIO(text))
    rows = []
    for row in reader:
        if any(cell.strip() for cell in row):
            rows.append(" | ".join(row))
    return "\n".join(rows) if rows else "[No data found in CSV]"


def _extract_pptx(raw_bytes: bytes) -> str:
    """Extract text from PPTX bytes using python-pptx."""
    from pptx import Presentation
    prs = Presentation(io.BytesIO(raw_bytes))
    all_text = []
    for i, slide in enumerate(prs.slides):
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_text.append(shape.text)
        if slide_text:
            all_text.append(f"--- Slide {i + 1} ---\n" + "\n".join(slide_text))
    return "\n\n".join(all_text) if all_text else "[No text found in presentation]"


def _extract_json(raw_bytes: bytes) -> str:
    """Pretty-print JSON content."""
    text = raw_bytes.decode("utf-8", errors="replace")
    try:
        data = json.loads(text)
        return json.dumps(data, indent=2, ensure_ascii=False)
    except json.JSONDecodeError:
        return text


def _extract_html(raw_bytes: bytes) -> str:
    """Extract text from HTML using BeautifulSoup."""
    from bs4 import BeautifulSoup
    text = raw_bytes.decode("utf-8", errors="replace")
    soup = BeautifulSoup(text, "html.parser")
    # Remove script and style elements
    for tag in soup(["script", "style"]):
        tag.decompose()
    return soup.get_text(separator="\n", strip=True)


def get_supported_extensions() -> list:
    """Return list of supported file extensions."""
    return [
        "pdf", "txt", "docx", "xlsx", "xls", "csv",
        "pptx", "json", "html", "htm", "md",
    ]


def get_file_icon(filename: str) -> str:
    """Return an emoji icon based on file type."""
    ext = filename.lower().rsplit(".", 1)[-1] if "." in filename else ""
    icons = {
        "pdf": "📄", "txt": "📝", "md": "📝",
        "docx": "📘", "doc": "📘",
        "xlsx": "📊", "xls": "📊", "csv": "📊",
        "pptx": "📙", "ppt": "📙",
        "json": "🔧", "html": "🌐", "htm": "🌐",
    }
    return icons.get(ext, "📎")
